import io
import time
import re
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from app.agents.state import AgentState
from app.agents.llm import get_model, extract_text_content
from app.core.storage import upload_file_artifact

logger = logging.getLogger("cortexflow")

PPT_PROMPT = """Create an outline for a professional 8-slide PowerPoint presentation on the topic:
Topic: {topic}

Output Rules:
- No conversational text or markdown intro.
- Follow strictly this format:

TITLE: Presentation Title
SUBTITLE: A compelling tagline

SLIDE:
Type: bullets
Title: Executive Summary
- Point one
- Point two
- Point three
- Point four

SLIDE:
Type: stats
Title: Key Metrics & Impact
- Market Size | $50B+
- Efficiency Gain | 45%
- Global Reach | 120+ Countries

SLIDE:
Type: conclusion
Title: Strategic Takeaways
- Key takeaway one
- Key takeaway two
- Key takeaway three
"""

def parse_ppt_response(content: str):
    res = {"title": "CortexFlow AI Presentation", "subtitle": "AI-Powered Deck", "slides": []}
    title_match = re.search(r"^TITLE:\s*(.+)", content, re.MULTILINE)
    if title_match:
        res["title"] = title_match.group(1).strip()
    
    sub_match = re.search(r"^SUBTITLE:\s*(.+)", content, re.MULTILINE)
    if sub_match:
        res["subtitle"] = sub_match.group(1).strip()
        
    slide_chunks = content.split("SLIDE:")[1:]
    for chunk in slide_chunks:
        lines = [l.strip() for l in chunk.strip().split("\n") if l.strip()]
        title = "Slide"
        slide_type = "bullets"
        items = []
        for line in lines:
            if line.lower().startswith("title:"):
                title = line.split(":", 1)[1].strip()
            elif line.lower().startswith("type:"):
                slide_type = line.split(":", 1)[1].strip().lower()
            elif line.startswith("-"):
                items.append(line.lstrip("-").strip())
        res["slides"].append({"title": title, "type": slide_type, "items": items})
    return res

async def ppt_node(state: AgentState) -> AgentState:
    prompt = state.get("prompt", "")
    llm = get_model("ppt")

    try:
        response = await llm.ainvoke(PPT_PROMPT.format(topic=prompt))
        parsed = parse_ppt_response(extract_text_content(response.content))

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

        # 1. Cover Slide
        cover_slide = prs.slides.add_slide(blank_layout)
        tb = cover_slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = parsed["title"]
        p.font.size = Pt(42)
        p.font.bold = True
        p.font.color.rgb = RGBColor(37, 99, 235)  # Blue
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = parsed["subtitle"]
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(100, 116, 139)
        p2.alignment = PP_ALIGN.CENTER

        # 2. Content Slides
        for slide_data in parsed["slides"]:
            slide = prs.slides.add_slide(blank_layout)
            
            # Slide Header
            header_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(1.0))
            htf = header_box.text_frame
            hp = htf.paragraphs[0]
            hp.text = slide_data["title"]
            hp.font.size = Pt(28)
            hp.font.bold = True
            hp.font.color.rgb = RGBColor(30, 41, 59)

            # Slide Content
            content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.5))
            ctf = content_box.text_frame
            ctf.word_wrap = True

            for idx, item in enumerate(slide_data["items"]):
                cp = ctf.paragraphs[0] if idx == 0 else ctf.add_paragraph()
                cp.text = f"•  {item}"
                cp.font.size = Pt(18)
                cp.font.color.rgb = RGBColor(51, 65, 85)
                cp.space_after = Pt(14)

        buffer = io.BytesIO()
        prs.save(buffer)
        ppt_bytes = buffer.getvalue()
        buffer.close()

        filename = f"presentation_{int(time.time())}.pptx"
        download_url = await upload_file_artifact(
            ppt_bytes, filename, "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

        return {
            **state,
            "response": f"""# 📊 Presentation Generated Successfully\n\n**Deck:** {parsed['title']}\n**Slides:** {len(parsed['slides']) + 1} Slides\n\n📥 **[Download PowerPoint Deck (.pptx)]({download_url})**\n\n*(Open in Microsoft PowerPoint, Google Slides, or Keynote)*""",
            "agent": "ppt"
        }
    except Exception as e:
        logger.error(f"PPT agent error: {e}")
        return {
            **state,
            "response": f"Failed to generate presentation: {str(e)}",
            "agent": "ppt"
        }
